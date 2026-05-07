"""Phase 3a — Extract text + page renders from PDF/PPTX originals.

For each row in inventory_classified.csv (target_kind != SKIP/UNMAPPED):
  - PDF:  pymupdf → content.md (page-by-page text), page_NNN.png (medium-res render)
  - PPTX: python-pptx → content.md (slide+notes text), page_NNN.png via render

Output: _migration/extracted/{kind}__{cat}__{slug}/
"""
from __future__ import annotations

import csv
import json
import shutil
import subprocess
from pathlib import Path

import fitz  # pymupdf
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
RAW_DIR = ROOT / "raw"
OUT_DIR = ROOT / "extracted"
CSV_IN = ROOT / "inventory_classified.csv"


def find_raw_file(page_id: str) -> Path | None:
    page_id_clean = page_id.replace("-", "")
    matches = sorted(RAW_DIR.glob(f"{page_id_clean}_*"))
    return matches[0] if matches else None


def extract_pdf(src: Path, dst: Path):
    doc = fitz.open(src)
    n_pages = len(doc)
    md_lines = [f"# {src.name}\n", f"_pages: {n_pages}_\n"]
    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        md_lines.append(f"\n## Page {i}\n\n{text}\n")
        pix = page.get_pixmap(matrix=fitz.Matrix(1.2, 1.2))
        pix.save(str(dst / f"page_{i:03d}.png"))
    (dst / "content.md").write_text("\n".join(md_lines), encoding="utf-8")
    doc.close()
    return n_pages


SOFFICE_PATHS = [
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    "/usr/local/bin/soffice",
    "/opt/homebrew/bin/soffice",
]


def pptx_to_pdf(src: Path, dst_pdf: Path) -> bool:
    """Convert .pptx → .pdf via LibreOffice (soffice). Returns True on success."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        for p in SOFFICE_PATHS:
            if Path(p).exists():
                soffice = p
                break
    if not soffice:
        return False
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(dst_pdf.parent), str(src)],
            capture_output=True, timeout=180, check=True,
        )
        # LibreOffice names output as src.stem + ".pdf"
        produced = dst_pdf.parent / (src.stem + ".pdf")
        if produced.exists() and produced != dst_pdf:
            produced.rename(dst_pdf)
        return dst_pdf.exists()
    except Exception:
        return False


def _walk_shapes(shapes):
    """Recursively walk pptx shapes (handles group shapes)."""
    for s in shapes:
        if getattr(s, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from _walk_shapes(s.shapes)
        else:
            yield s


def extract_pptx(src: Path, dst: Path) -> int:
    prs = Presentation(src)

    # 1) Speaker notes (only thing python-pptx reliably gets for infographic-style decks)
    notes_blocks: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
            if text:
                notes_blocks.append(f"## Slide {i} Notes\n\n{text}\n")

    # 2) Convert to PDF via LibreOffice → text + page renders
    tmp_pdf = dst / "_render.pdf"
    if not pptx_to_pdf(src, tmp_pdf):
        # Fallback: python-pptx text-only walk (rare)
        md_lines = [f"# {src.name}\n", f"_slides: {len(prs.slides)}_\n"]
        for i, slide in enumerate(prs.slides, 1):
            md_lines.append(f"\n## Slide {i}\n")
            for shape in _walk_shapes(slide.shapes):
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            md_lines.append(f"- {para.text.strip()}")
        (dst / "content.md").write_text("\n".join(md_lines), encoding="utf-8")
        (dst / "RENDER_SKIPPED.txt").write_text(
            "LibreOffice not installed — pptx rendered as text only.",
            encoding="utf-8",
        )
        return len(prs.slides)

    doc = fitz.open(tmp_pdf)
    n_pages = len(doc)
    md_lines = [f"# {src.name}", f"_slides: {n_pages}_  (rendered via LibreOffice)\n"]
    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        md_lines.append(f"\n## Slide {i}\n\n{text}\n")
        pix = page.get_pixmap(matrix=fitz.Matrix(1.2, 1.2))
        pix.save(str(dst / f"page_{i:03d}.png"))
    if notes_blocks:
        md_lines.append("\n---\n\n# Speaker Notes\n")
        md_lines.extend(notes_blocks)
    (dst / "content.md").write_text("\n".join(md_lines), encoding="utf-8")
    doc.close()
    tmp_pdf.unlink()
    return n_pages


def main():
    rows = list(csv.DictReader(CSV_IN.open(encoding="utf-8")))
    targets = [r for r in rows if r["target_kind"] in ("decks", "handouts", "lab-reports")]
    print(f"Extracting from {len(targets)} files…")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    summary = {"ok": 0, "fail": 0, "no_render": 0}

    for r in targets:
        kind, cat, slug = r["target_kind"], r["category_dir"], r["slug"]
        dst_name = f"{kind}__{cat}__{slug}"
        dst = OUT_DIR / dst_name
        dst.mkdir(parents=True, exist_ok=True)

        src = find_raw_file(r["page_id"])
        if not src:
            print(f"  ✗ {dst_name}: raw not found (page_id={r['page_id']})")
            summary["fail"] += 1
            continue

        try:
            if src.suffix.lower() == ".pdf":
                pages = extract_pdf(src, dst)
            elif src.suffix.lower() == ".pptx":
                pages = extract_pptx(src, dst)
            else:
                print(f"  ? {dst_name}: unsupported {src.suffix}")
                summary["fail"] += 1
                continue
            (dst / "metadata.json").write_text(json.dumps({
                "page_id": r["page_id"],
                "title": r["title"],
                "category": r["category"],
                "audience": r["audience"],
                "detail": r["detail"],
                "target_kind": kind, "category_dir": cat, "slug": slug,
                "source_file": src.name,
                "page_count": pages,
            }, ensure_ascii=False, indent=2), encoding="utf-8")
            print(f"  ✓ {dst_name} ({pages}p)")
            summary["ok"] += 1
            if (dst / "RENDER_SKIPPED.txt").exists():
                summary["no_render"] += 1
        except Exception as e:
            print(f"  ✗ {dst_name}: {e}")
            summary["fail"] += 1

    print(f"\nDone: {summary}")
    if summary["no_render"]:
        print(f"⚠️ {summary['no_render']} PPTX rendered as text only (LibreOffice missing).")


if __name__ == "__main__":
    main()
